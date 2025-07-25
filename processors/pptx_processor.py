import logging
from pptx import Presentation
from config.settings import settings
from utils.chunking import chunk_text

CHUNK_SIZE = settings.CHUNK_SIZE
CHUNK_OVERLAP = settings.CHUNK_OVERLAP

logger = logging.getLogger(__name__)

class PPTXProcessor:
    """
    A processor for extracting text from PPTX documents and chunking it.
    """
    def __init__(self):
        """
        Initializes the PPTXProcessor.
        """
        logger.info("PPTXProcessor initialized.")

    def process(self, file_path: str) -> list[str]:
        """
        Extracts text from a PPTX file, chunks it, and returns a list of text chunks.

        Args:
            file_path (str): The path to the PPTX file.

        Returns:
            list[str]: A list of text chunks extracted from the PPTX.
        """
        try:
            presentation = Presentation(file_path)
            full_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):  # Reason: Check if the shape has text content.
                        full_text.append(shape.text)
            
            combined_text = "\n".join(full_text)
            chunks = chunk_text(combined_text, CHUNK_SIZE, CHUNK_OVERLAP)
            logger.info(f"Successfully processed PPTX file: {file_path}. Extracted {len(chunks)} chunks.")
            return chunks
        except Exception as e:
            logger.error(f"Error processing PPTX file {file_path}: {e}")
            return []